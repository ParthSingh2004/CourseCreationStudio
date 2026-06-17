import fitz  # PyMuPDF
import docx
import io
from pptx import Presentation

def extract_text_from_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text

def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text(file_bytes: bytes, filename: str) -> str:
    if filename.lower().endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif filename.lower().endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif filename.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif filename.lower().endswith(".txt"):
        return file_bytes.decode('utf-8', errors='ignore')
    return ""
